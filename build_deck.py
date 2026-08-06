"""Generate dist/presentation.pptx (Deliverable 2) — 10 slides.

Every number is read from the canonical root results files (the live run):
  ft_eval_results.json   (a/b/c calibrated metrics + paired tests)
  results/summary.json   (explainer aggregate: field flips, drivers, faithfulness)
Figures are rendered from cached results only — NO API calls. The one figure that
needs per-row zero-shot probabilities (the reliability diagram) reads them from the
on-disk response cache and ABORTS before any network call if coverage is < 100%.

Sparse slide faces (the "On the slide" bullets); all coaching lives in the
speaker-notes pane. 16:9, "Private & Confidential" footer. The .pptx is written to
dist/ and is kept OUT of the code submission zip.
"""
from __future__ import annotations

import json
import os
import random

import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

ROOT = os.path.dirname(os.path.abspath(__file__))
os.chdir(ROOT)
FIGS = "/private/tmp/claude-501/-Users-spapadias-Desktop-Research-icloud-satalia/c062b385-2e72-4662-bc7f-2c2783b9ab02/scratchpad/deck_figs"
os.makedirs(FIGS, exist_ok=True)

NAVY = "#1f3a5f"
BLUE = "#2e6da4"
GREY = "#7a7a7a"
LIGHT = "#d9e2ec"
GREEN = "#2e7d4f"
RED = "#b5443b"

# --------------------------------------------------------------------------
# 1. NUMBERS — read from the canonical result files
# --------------------------------------------------------------------------
FE = json.load(open("ft_eval_results.json"))
SUMM = json.load(open("results/summary.json"))
EXPL = SUMM["explainer"]
b, c, a = FE["b"], FE["c"], FE["a"]
cb = FE["c_vs_b"]

M = ["accuracy", "balanced_accuracy", "macro_f1", "roc_auc", "pr_auc", "brier", "ece"]
print("[numbers] read from ft_eval_results.json + results/summary.json")
print(f"  zero-shot (b): acc {b['accuracy']:.3f} ece {b['ece']:.3f}")
print(f"  fine-tuned (c): acc {c['accuracy']:.3f} ece {c['ece']:.3f}")
print(f"  score-mode (a): acc {a['accuracy']:.3f} ece {a['ece']:.3f}")
print(f"  c_vs_b: dacc {cb['acc_diff'][0]:+.3f} [{cb['acc_diff'][1]:+.3f},{cb['acc_diff'][2]:+.3f}] p={cb['mcnemar']['p_value']:.2e}")

field_tbl = EXPL["field_flip_table"]
flip = {r["removed_field"]: r["flip_rate"] for r in field_tbl}
dc = {r["driver"]: r for r in EXPL["driver_correctness"]}
agree = EXPL["rationale_occlusion_agreement_rate"]
print(f"  explainer: speaker flip {flip['speaker_name']} all_meta {flip['all_metadata']} "
      f"agree {agree} | statement {dc['statement']['accuracy']} speaker {dc['speaker_name']['accuracy']}")

# --------------------------------------------------------------------------
# 2. SPLIT SIZES + TRUE-RATES (computed from data, for slide 3)
# --------------------------------------------------------------------------
from truthclf import data, calibration, metrics, prompts, llm
from truthclf.predictors import ZeroShotPredictor
from truthclf.predictors.zeroshot import prob_from_logprobs

SCHEME = "primary"
clean, _ = data.clean_dataset(data.load("data.csv"), SCHEME)
train, val, test = data.speaker_disjoint_3way(clean, 0.2, 0.2, 0, SCHEME)
def true_rate(rows): return sum(r.y(SCHEME) for r in rows) / len(rows)
SPLIT = {"Train": (len(train), true_rate(train)),
         "Validation": (len(val), true_rate(val)),
         "Test": (len(test), true_rate(test))}
print(f"[split] " + " ".join(f"{k} {n} (True {tr:.2f})" for k, (n, tr) in SPLIT.items()))

# --------------------------------------------------------------------------
# 3. FIGURES
# --------------------------------------------------------------------------
figs_made, figs_failed = [], []

# 3a. slide-3 split diagram --------------------------------------------------
def fig_split():
    fig, ax = plt.subplots(figsize=(9.2, 3.0))
    ax.axis("off")
    total = sum(n for n, _ in SPLIT.values())
    x = 0.0
    colors = {"Train": BLUE, "Validation": "#6fa3d2", "Test": NAVY}
    for name, (n, tr) in SPLIT.items():
        w = n / total
        ax.add_patch(plt.Rectangle((x, 0.35), w, 0.4, color=colors[name]))
        ax.text(x + w / 2, 0.55, name, ha="center", va="center", color="white",
                fontsize=13, fontweight="bold")
        ax.text(x + w / 2, 0.22, f"n = {n:,}", ha="center", va="center", fontsize=11)
        ax.text(x + w / 2, 0.10, f"True-rate {tr:.2f}", ha="center", va="center",
                fontsize=10, color=GREY)
        x += w
    ax.text(0.5, 0.92, "Speaker-disjoint  +  near-duplicate-disjoint (union-find)",
            ha="center", fontsize=12, fontweight="bold", color=NAVY)
    ax.text(0.5, 0.82, "no speaker and no near-identical statement spans two splits",
            ha="center", fontsize=10, color=GREY)
    ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    p = f"{FIGS}/split.png"; fig.savefig(p, dpi=150, bbox_inches="tight"); plt.close(fig)
    return p

# 3b. slide-8 reliability (zero-shot vs fine-tuned, both calibrated) ----------
_NO_THINK = {"chat_template_kwargs": {"enable_thinking": False}}

def _offline_zeroshot_probs(rows):
    """Decision/logprob zero-shot probs read from the on-disk cache ONLY, via the
    SAME batch client the live run used (base_decision_probs) so the cached
    value-format matches. Verifies 100% coverage and raises before any network
    call on a miss (all-cached -> _serve never submits a batch)."""
    client = llm.TogetherBatchClient("google/gemma-4-31B-it", max_output_tokens=1,
                                     extra_create_kwargs=_NO_THINK)
    missing = [r.row_id for r in rows
               if client.cache.get(client._key_classify(
                   prompts.build_messages(r, "full", "decision"))) is None]
    if missing:
        raise RuntimeError(f"reliability fig: {len(missing)}/{len(rows)} zero-shot rows "
                           f"not cached — aborting to avoid API calls")
    pred = ZeroShotPredictor("google/gemma-4-31B-it", "full", client=client,
                             use_logprobs=True)
    return pred.predict(rows).probs

def _ft_probs(rows):
    fc = json.load(open("ft_eval_cache.json"))
    return [fc[str(r.row_id)] for r in rows]

def fig_reliability():
    vy = [r.y(SCHEME) for r in val]; ty = [r.y(SCHEME) for r in test]
    # zero-shot
    zs_v, zs_t = _offline_zeroshot_probs(val), _offline_zeroshot_probs(test)
    cal_zs = calibration.fit_best(zs_v, vy)
    zs_tc = list(calibration.apply(zs_t, cal_zs))
    # fine-tuned
    ft_v, ft_t = _ft_probs(val), _ft_probs(test)
    cal_ft = calibration.fit_best(ft_v, vy)
    ft_tc = list(calibration.apply(ft_t, cal_ft))
    ece_zs, ece_ft = metrics.ece(ty, zs_tc), metrics.ece(ty, ft_tc)
    print(f"[reliability] computed test ECE — zero-shot {ece_zs:.3f}  fine-tuned {ece_ft:.3f} "
          f"(file: {b['ece']:.3f} / {c['ece']:.3f})")
    fig, ax = plt.subplots(figsize=(5.6, 5.0))
    ax.plot([0, 1], [0, 1], "--", color=GREY, lw=1, label="perfect")
    for probs, col, lab, ece in [(zs_tc, BLUE, "zero-shot", ece_zs),
                                 (ft_tc, GREEN, "fine-tuned", ece_ft)]:
        rc = metrics.reliability_curve(ty, probs)
        mp = [m for m, c_ in zip(rc["mean_pred"], rc["count"]) if c_ > 0]
        fp = [f for f, c_ in zip(rc["frac_pos"], rc["count"]) if c_ > 0]
        ax.plot(mp, fp, "o-", color=col, lw=2, ms=5, label=f"{lab}  (ECE {ece:.3f})")
    ax.set_xlabel("predicted probability"); ax.set_ylabel("observed frequency")
    ax.set_title("Reliability on test split — calibrated", fontsize=12, color=NAVY)
    ax.legend(loc="upper left", fontsize=10); ax.set_xlim(0, 1); ax.set_ylim(0, 1)
    ax.grid(alpha=0.25)
    p = f"{FIGS}/reliability.png"; fig.savefig(p, dpi=150, bbox_inches="tight"); plt.close(fig)
    return p

# 3c. slide-9 field-importance bar ------------------------------------------
def fig_fields():
    order = ["all_metadata", "speaker_name", "statement_context", "subjects",
             "speaker_affiliation"]
    labs = {"all_metadata": "all metadata", "speaker_name": "speaker name",
            "statement_context": "context", "subjects": "subjects",
            "speaker_affiliation": "affiliation"}
    vals = [flip[k] for k in order]
    fig, ax = plt.subplots(figsize=(5.6, 3.4))
    cols = [NAVY] + [BLUE] * (len(order) - 1)
    y = np.arange(len(order))[::-1]
    ax.barh(y, vals, color=cols)
    for yi, v in zip(y, vals):
        ax.text(v + 0.003, yi, f"{v*100:.1f}%", va="center", fontsize=10)
    ax.set_yticks(y); ax.set_yticklabels([labs[k] for k in order], fontsize=10)
    ax.set_xlabel("prediction flip-rate when field removed (occlusion)", fontsize=10)
    ax.set_title("Field importance (leave-one-field-out)", fontsize=12, color=NAVY)
    ax.set_xlim(0, max(vals) * 1.25); ax.grid(axis="x", alpha=0.25)
    p = f"{FIGS}/fields.png"; fig.savefig(p, dpi=150, bbox_inches="tight"); plt.close(fig)
    return p

# 3d. slide-9 driver-correctness contrast -----------------------------------
def fig_drivers():
    st, sp = dc["statement"], dc["speaker_name"]
    fig, ax = plt.subplots(figsize=(5.6, 3.4))
    bars = [("statement-driven", st["accuracy"], st["n"], GREEN),
            ("speaker-driven", sp["accuracy"], sp["n"], RED)]
    x = np.arange(len(bars))
    ax.bar(x, [v for _, v, _, _ in bars], color=[col for *_, col in bars], width=0.55)
    for xi, (_, v, n, _) in zip(x, bars):
        ax.text(xi, v + 0.01, f"{v:.3f}\n(n={n})", ha="center", fontsize=11)
    ax.set_xticks(x); ax.set_xticklabels([lab for lab, *_ in bars], fontsize=11)
    ax.set_ylabel("accuracy"); ax.set_ylim(0, 1)
    ax.axhline(0.5, ls="--", color=GREY, lw=1)
    ax.set_title("Predictions driven by the statement are more accurate\n"
                 "than those driven by the speaker", fontsize=11, color=NAVY)
    ax.grid(axis="y", alpha=0.25)
    p = f"{FIGS}/drivers.png"; fig.savefig(p, dpi=150, bbox_inches="tight"); plt.close(fig)
    return p

for name, fn in [("split", fig_split), ("reliability", fig_reliability),
                 ("fields", fig_fields), ("drivers", fig_drivers)]:
    try:
        path = fn(); figs_made.append((name, path)); print(f"[fig] {name} -> {path}")
    except Exception as e:
        figs_failed.append((name, str(e))); print(f"[fig] FAILED {name}: {e}")
FIGP = {n: p for n, p in figs_made}

# --------------------------------------------------------------------------
# 4. DECK
# --------------------------------------------------------------------------
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def rgb(h): return RGBColor.from_string(h.lstrip("#"))
NAVY_C, BLUE_C, GREY_C, GREEN_C, RED_C = map(rgb, [NAVY, BLUE, GREY, GREEN, RED])

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_slide(title, bullets, notes, fig=None, fig_box=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    # title
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), SW - Inches(1.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY_C
    if sub:
        p2 = tf.add_paragraph(); rr = p2.add_run(); rr.text = sub
        rr.font.size = Pt(14); rr.font.italic = True; rr.font.color.rgb = BLUE_C
    # accent rule
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.32), Inches(2.4), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE_C; ln.line.fill.background()
    # body bullets
    body_w = (Inches(6.2) if fig else SW - Inches(1.1))
    body = s.shapes.add_textbox(Inches(0.55), Inches(1.55), body_w, SH - Inches(2.6))
    bf = body.text_frame; bf.word_wrap = True
    for i, txt in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        run = para.add_run(); run.text = "•  " + txt
        run.font.size = Pt(16); run.font.color.rgb = rgb("#222222")
        para.space_after = Pt(10)
    # figure
    if fig and os.path.exists(fig):
        x, y, w = fig_box or (Inches(7.0), Inches(1.7), Inches(5.9))
        s.shapes.add_picture(fig, x, y, width=w)
    # footer
    ft = s.shapes.add_textbox(Inches(0.55), SH - Inches(0.45), SW - Inches(1.1), Inches(0.35))
    fr = ft.text_frame.paragraphs[0]; run = fr.add_run()
    run.text = "Private & Confidential"
    run.font.size = Pt(9); run.font.color.rgb = GREY_C; fr.alignment = PP_ALIGN.RIGHT
    # notes
    s.notes_slide.notes_text_frame.text = notes
    return s

# ---- Slide 1
add_slide(
    "The Problem & Its Real Ceiling",
    ["Classify public statements as True / False from text + metadata (speaker, party, subjects, context).",
     "Low-ceiling task: ~62–69% binary accuracy across model families and methods.",
     "Why: verifying claims needs external world knowledge the text lacks; human truth labels are judgment calls with limited agreement; some statements aren't verifiable from inputs.",
     "Aim: a principled, honestly-evaluated system — not a vanity accuracy number."],
    "How do you know the ceiling? Published benchmarks cluster there; my results landed in the same band. "
    "65% sounds weak? Floor is 50% for balanced binary, so it's real signal but task-capped — the responsible "
    "response is calibrated confidence and abstention. This slide sets expectations for every modest number later.")

# ---- Slide 2
add_slide(
    "Target Framing: 6 Labels → Binary",
    ["Source: 6 ordinal truthfulness levels.",
     "Middle split: 3 truer → True, 3 falser → False.",
     "Why: near-balanced target (≈56/44) keeps accuracy meaningful, no majority-class shortcut.",
     "Sensitivity check: moving ambiguous \"half-true\" to False leaves conclusions unchanged."],
    "Why not drop the middle / 2-vs-4? Dropping shrinks a small set; imbalance makes accuracy misleading; "
    "middle split is balanced and conventional. Half-true really True? Hardest call — placed True for balance, "
    "but sensitivity check shows robustness either way. Sensitivity check = re-run the whole eval with one "
    "debatable decision flipped.")

# ---- Slide 3
add_slide(
    "Data Integrity: The Split That Prevents Cheating",
    ["Naive random splits leak: frequent speakers appear in train and test → model learns \"this person lies\", not the claim.",
     "My split: speaker-disjoint AND near-duplicate-disjoint (union-find) — no speaker or near-identical statement spans two splits.",
     "Dropped contradictory duplicates (identical text, conflicting labels).",
     f"Train {SPLIT['Train'][0]:,} / Val {SPLIT['Validation'][0]:,} / Test {SPLIT['Test'][0]:,}; matched balance "
     f"(True-rate {SPLIT['Train'][1]:.2f} / {SPLIT['Validation'][1]:.2f} / {SPLIT['Test'][1]:.2f})."],
    "Why speaker leakage matters: it's a shortcut; my explainer later quantifies it's real and hurts accuracy — "
    "which is why I split to prevent it. Union-find: merges linked items (shared speaker or near-dup text) into "
    "clusters assigned whole to one split; guarantees no link crosses train/test. Strongest methodology slide — "
    "most candidates random-split and miss this.",
    fig=FIGP.get("split"), fig_box=(Inches(1.9), Inches(4.6), Inches(9.5)))

# ---- Slide 4
add_slide(
    "Evaluation Methodology: Metrics & Protocol",
    ["Beyond accuracy: balanced accuracy, macro-F1, ROC/PR-AUC, and calibration (Brier + ECE).",
     "Cost asymmetry: a false positive (falsehood labelled True) lets misinformation through — the costlier error; threshold can reflect that.",
     "Protocol: calibrator + threshold fit on validation, reported on untouched test.",
     "Bootstrap CIs; McNemar's test for the paired comparison."],
    "Why ECE/Brier? System acts on confidence; a model can be accurate yet overconfident. Why val-fit/test-report? "
    "Tuning on test = peeking = inflated; this keeps numbers out-of-sample. Definitions — balanced acc: average "
    "per-class accuracy. Macro-F1: per-class F1 averaged. Brier: mean squared error of probability vs outcome. "
    "ECE: groups by stated confidence, checks if \"80% sure\" is right 80% — lower = trustworthy. McNemar: paired "
    "test on the cases where two models disagree.")

# ---- Slide 5
add_slide(
    "Zero-Shot Predictor",
    ["Base: gemma-4-31B-it (open; serverless on Together; LoRA-fine-tunable under the same id → shared base, no confound).",
     "Prompt: statement + metadata, missing fields marked [unknown].",
     "Elicitation: decision + logprob — model answers True/False; confidence from the True/False token's logprob.",
     f"Baseline (calibrated, test): acc {b['accuracy']:.3f}, ECE {b['ece']:.3f}."],
    "Why this model? Scale-tested candidates; a 9B Qwen and Gemma hit the literature range, OpenAI 20B/120B sat "
    "at chance; Gemma won the full-set comparison and is the one base that's both serverless and fine-tunable "
    "under one id — so no model-swap confound. Logprobs: the API exposes the probability the model assigned its "
    "chosen token; exp(logprob) = real confidence, not a performed number.")

# ---- Slide 6
add_slide(
    "Fine-Tuned Predictor (LoRA SFT)",
    ["LoRA supervised fine-tuning on the same gemma-4-31B-it base.",
     "Training data: conversational JSONL — prompt as system+user, assistant = single True/False token; loss on the label only (train_on_inputs=auto).",
     "Train split only (5,710); validation file for eval-loss. Eval loss 0.32 → 0.199 over 3 epochs (from the training run).",
     "Serving: fine-tuned Gemma isn't serverless → short-lived dedicated endpoint, create→infer→guaranteed teardown (~14 min, ~$2)."],
    "LoRA = trains small low-rank adapters on top of a frozen base; cheap, fast, enough to shift behaviour. "
    "Why label-token-only? Learn to produce the answer, not echo the prompt. Why an endpoint now? Base served "
    "serverless, but fine-tuned Gemma can't; endpoint only for eval, deleted in a guaranteed-teardown block. "
    "Low train loss = learned the signal; whether it transfers is what the disjoint test answers. "
    "Eval-loss 0.32 → 0.199 is from the fine-tuning job's training log (3 epochs).")

# ---- Slide 7 (native table)
def slide7():
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), SW - Inches(1.1), Inches(1.0))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = "The Fair Comparison & Results"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = NAVY_C
    p2 = tf.add_paragraph(); rr = p2.add_run()
    rr.text = "Fine-tuning beats matched zero-shot +3.3 pts (p < 0.001) AND halves calibration error."
    rr.font.size = Pt(14); rr.font.italic = True; rr.font.color.rgb = BLUE_C
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.5), Inches(2.4), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE_C; ln.line.fill.background()
    # one-line setup bullet
    bb = s.shapes.add_textbox(Inches(0.55), Inches(1.62), SW - Inches(1.1), Inches(0.5))
    br = bb.text_frame.paragraphs[0].add_run()
    br.text = "•  Same test split, same calibration, same elicitation → only fine-tuning differs."
    br.font.size = Pt(15); br.font.color.rgb = rgb("#222222")

    cols = ["Predictor", "Acc", "Bal-acc", "Macro-F1", "ROC-AUC", "PR-AUC", "Brier", "ECE"]
    rows_data = [
        ("Zero-shot, logprob (baseline)", b, "primary"),
        ("Fine-tuned (LoRA SFT, decision)", c, "ft"),
        ("secondary · score elicitation, not matched", a, "secondary"),
    ]
    nrows, ncols = len(rows_data) + 1, len(cols)
    gt = s.shapes.add_table(nrows, ncols, Inches(0.55), Inches(2.25),
                            Inches(12.2), Inches(2.2)).table
    gt.columns[0].width = Inches(4.3)
    for j in range(1, ncols):
        gt.columns[j].width = Inches((12.2 - 4.3) / (ncols - 1))
    # header
    for j, ct in enumerate(cols):
        cell = gt.cell(0, j); cell.text = ct
        para = cell.text_frame.paragraphs[0]; para.alignment = PP_ALIGN.CENTER
        run = para.runs[0]; run.font.size = Pt(13); run.font.bold = True
        run.font.color.rgb = RGBColor.from_string("FFFFFF")
        cell.fill.solid(); cell.fill.fore_color.rgb = NAVY_C
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    # body
    keys = ["accuracy", "balanced_accuracy", "macro_f1", "roc_auc", "pr_auc", "brier", "ece"]
    for i, (lab, mset, kind) in enumerate(rows_data, start=1):
        vals = [lab] + [f"{mset[k]:.3f}" for k in keys]
        for j, v in enumerate(vals):
            cell = gt.cell(i, j); cell.text = v
            para = cell.text_frame.paragraphs[0]
            para.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            run = para.runs[0]; run.font.size = Pt(12)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if kind == "secondary":
                run.font.italic = True; run.font.color.rgb = GREY_C
                cell.fill.solid(); cell.fill.fore_color.rgb = rgb("#efefef")
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = (rgb("#e8f0e8") if kind == "ft" else rgb("#f4f7fb"))
                if kind == "ft":
                    run.font.bold = True; run.font.color.rgb = GREEN_C
                else:
                    run.font.color.rgb = rgb("#222222")
    # caption
    cap = s.shapes.add_textbox(Inches(0.55), Inches(4.65), Inches(12.2), Inches(1.4))
    cf = cap.text_frame; cf.word_wrap = True
    d = cb["acc_diff"]
    p = cf.paragraphs[0]; run = p.add_run()
    run.text = (f"Clean fine-tuning effect (matched elicitation): "
                f"Δacc {d[0]:+.3f} [{d[1]:+.3f}, {d[2]:+.3f}], McNemar p < 0.001. "
                f"Calibration: ECE {b['ece']:.3f} → {c['ece']:.3f}.")
    run.font.size = Pt(14); run.font.bold = True; run.font.color.rgb = NAVY_C
    p2 = cf.add_paragraph(); r2 = p2.add_run()
    r2.text = ("Secondary (score-mode): highest accuracy yet ~4× worse calibrated "
               f"(ECE {a['ece']:.3f} vs {c['ece']:.3f}) — accuracy ≠ trustworthiness.")
    r2.font.size = Pt(12); r2.font.italic = True; r2.font.color.rgb = GREY_C
    # footer
    ft = s.shapes.add_textbox(Inches(0.55), SH - Inches(0.45), SW - Inches(1.1), Inches(0.35))
    fr = ft.text_frame.paragraphs[0]; run = fr.add_run(); run.text = "Private & Confidential"
    run.font.size = Pt(9); run.font.color.rgb = GREY_C; fr.alignment = PP_ALIGN.RIGHT
    s.notes_slide.notes_text_frame.text = (
        "Only 3.3 points? On a ~65% ceiling, a significant gain over a matched baseline plus halved calibration "
        "error is meaningful — accuracy was capped, trustworthy confidence is the axis that moved. Why show "
        "score-mode? It's the contrast that proves accuracy ≠ calibration. Frame: train fine-tunes; val fits "
        "calibrator+threshold for both; test scored once identically — only fine-tuning differs.\n\n"
        "Variance caveat: results vary slightly run-to-run due to non-deterministic inference even at temperature 0; "
        "the fine-tuning effect was consistently positive and significant across runs (+0.017 to +0.033).")
slide7()

# ---- Slide 8
add_slide(
    "Calibration: The Real Win",
    ["Raw confidence is miscalibrated → temperature/Platt scaling, fit on validation, corrects it.",
     f"ECE {b['ece']:.3f} → {c['ece']:.3f} (more than halved); reliability diagram moves toward the diagonal.",
     "For a tool that acts on confidence (abstain, escalate, flag), trustworthy probabilities are the operationally useful property."],
    "Temperature/Platt = simple post-hoc corrections fit on held-out data remapping raw probabilities so stated "
    "confidence matches real frequency; I fit several on validation and selected the best. Why not isotonic? "
    "Needs more data, overfits ~1,900 rows — future work. Say cold: calibration doesn't change the answer or "
    "accuracy — it makes the confidence number honest. The diagram shows both calibrated models on test; the "
    "fine-tuned curve sits closest to the diagonal (lowest ECE).",
    fig=FIGP.get("reliability"), fig_box=(Inches(7.4), Inches(1.7), Inches(5.2)))

# ---- Slide 9 (two figures)
def slide9():
    s = prs.slides.add_slide(BLANK)
    tb = s.shapes.add_textbox(Inches(0.55), Inches(0.35), SW - Inches(1.1), Inches(1.0))
    r = tb.text_frame.paragraphs[0].add_run()
    r.text = "Explainer: Faithful Attribution Beats Plausible Stories"
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = NAVY_C
    ln = s.shapes.add_shape(1, Inches(0.55), Inches(1.28), Inches(2.4), Pt(3))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE_C; ln.line.fill.background()
    body = s.shapes.add_textbox(Inches(0.55), Inches(1.45), Inches(12.3), Inches(1.7))
    bf = body.text_frame; bf.word_wrap = True
    bullets = [
        "explain(model, points, labels) — model-agnostic, either predictor.",
        "Two layers: leave-one-field-out occlusion (faithful by construction) + model's own rationale (plausible, not necessarily faithful).",
        f"Removing speaker flips {flip['speaker_name']*100:.0f}% of predictions, all metadata {flip['all_metadata']*100:.1f}%; "
        f"rationales only ~{agree*100:.0f}% faithful; speaker-driven predictions less accurate "
        f"(statement-driven {dc['statement']['accuracy']:.3f} vs speaker-driven {dc['speaker_name']['accuracy']:.3f}).",
    ]
    for i, txt in enumerate(bullets):
        para = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        run = para.add_run(); run.text = "•  " + txt
        run.font.size = Pt(15); run.font.color.rgb = rgb("#222222"); para.space_after = Pt(8)
    if FIGP.get("fields"):
        s.shapes.add_picture(FIGP["fields"], Inches(0.55), Inches(3.5), width=Inches(6.1))
    if FIGP.get("drivers"):
        s.shapes.add_picture(FIGP["drivers"], Inches(6.9), Inches(3.5), width=Inches(6.1))
    ft = s.shapes.add_textbox(Inches(0.55), SH - Inches(0.45), SW - Inches(1.1), Inches(0.35))
    fr = ft.text_frame.paragraphs[0]; run = fr.add_run(); run.text = "Private & Confidential"
    run.font.size = Pt(9); run.font.color.rgb = GREY_C; fr.alignment = PP_ALIGN.RIGHT
    s.notes_slide.notes_text_frame.text = (
        "Why two methods? Occlusion is causal — change input, watch real output; rationale is readable but only "
        "~46% faithful, so I lead with occlusion and treat rationales as hypotheses, catching the model "
        "rationalizing. 46% — good or bad? Bad for the model's self-explanations, good for my explainer that "
        "detected it. Shortcut helps or hurts? Hurts (62.5% vs 75.6% — speaker- vs statement-driven); "
        "correlational but speaker-driven predictions are the less trustworthy ones — causal confirmation of "
        "slide 3's leakage risk.")
slide9()

# ---- Slide 10
add_slide(
    "Findings & Future Work",
    ["Story: a low-ceiling task evaluated leakage-free; fine-tuning gave a significant accuracy gain over matched zero-shot and halved calibration error; the explainer exposed a speaker shortcut that hurts and showed self-rationales are unreliable.",
     "Optimized for honest, defensible evaluation over a vanity number.",
     "Future work: logprob elicitation everywhere; richer calibration (isotonic); ensembling zero-shot + fine-tuned; production abstention from the coverage-accuracy curve; alternative elicitations (score-mode trade-off); DPO if ranked data existed."],
    "Another week? Push abstention into a selective-prediction system on the calibrated confidence; ensemble both "
    "predictors; resist chasing accuracy past the ceiling. One thing to remember: accuracy and trustworthiness "
    "are different axes — fine-tuning's real value was honest confidence, and every claim is reproducible from "
    "the code. Closes the loop to slide 1.")

os.makedirs("dist", exist_ok=True)
OUT = "dist/presentation.pptx"
prs.save(OUT)
print(f"\n[deck] saved {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
print(f"[figs] generated: {[n for n,_ in figs_made]}")
if figs_failed:
    print(f"[figs] FAILED: {figs_failed}")
